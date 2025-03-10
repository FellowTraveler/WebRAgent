import os
import re
import uuid
import mimetypes
import logging
import tempfile
import shutil
import zipfile
from pathlib import Path
from typing import Optional, Dict, Any, List, Tuple

# Docling imports for advanced document processing
from docling.document_converter import DocumentConverter
from docling.datamodel.base_models import InputFormat
from docling import document_converter

# Keep original imports for fallback handling
from PyPDF2 import PdfReader
from bs4 import BeautifulSoup
import markdown
from werkzeug.datastructures import FileStorage

from app.models.document import Document
from app.models.collection import Collection
from app.services.qdrant_service import QdrantService

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class DocumentService:
    """Service for processing and managing documents"""
    
    def __init__(self):
        """Initialize document service"""
        self.upload_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), 
                                       'data', 'uploads')
        os.makedirs(self.upload_dir, exist_ok=True)
        self.qdrant_service = QdrantService()
        
        # Initialize docling document converter with allowed formats
        self.converter = DocumentConverter(
            allowed_formats=[
                InputFormat.PDF,
                InputFormat.DOCX,
                InputFormat.PPTX,
                InputFormat.XLSX,
                InputFormat.HTML,
                InputFormat.MD,
                InputFormat.ASCIIDOC,
                InputFormat.CSV,
                InputFormat.XML_JATS,
                InputFormat.IMAGE,
            ]
        )
        
        # Extension to format mapping for supported types
        self.extension_map = {
            # PDF files
            '.pdf': InputFormat.PDF,
            
            # Office documents
            '.doc': InputFormat.DOCX,
            '.docx': InputFormat.DOCX,
            '.ppt': InputFormat.PPTX, 
            '.pptx': InputFormat.PPTX,
            '.xls': InputFormat.XLSX,
            '.xlsx': InputFormat.XLSX,
            
            # Web/markup formats
            '.html': InputFormat.HTML,
            '.htm': InputFormat.HTML,
            '.md': InputFormat.MD,
            '.adoc': InputFormat.ASCIIDOC,
            
            # Data formats
            '.csv': InputFormat.CSV,
            '.xml': InputFormat.XML_JATS,
            
            # Image files (with OCR)
            '.jpg': InputFormat.IMAGE,
            '.jpeg': InputFormat.IMAGE,
            '.png': InputFormat.IMAGE,
            '.tif': InputFormat.IMAGE,
            '.tiff': InputFormat.IMAGE,
            
            # Archive formats - special handling
            '.zip': 'ZIP_ARCHIVE',  # Not an InputFormat, handled specially
        }
        
    def save_uploaded_file(self, file):
        """
        Save uploaded file to disk
        
        Args:
            file: The uploaded file object
            
        Returns:
            str: Path to the saved file
        """
        filename = str(uuid.uuid4()) + "_" + file.filename
        file_path = os.path.join(self.upload_dir, filename)
        file.save(file_path)
        return file_path
        
    def process_zip_file(self, zip_file_path):
        """
        Extract files from a ZIP archive and return paths to extracted files
        
        Args:
            zip_file_path: Path to the ZIP file
            
        Returns:
            List[Tuple[str, FileStorage]]: List of (extracted file path, file object) tuples
        """
        extracted_files = []
        extract_dir = os.path.join(self.upload_dir, str(uuid.uuid4()) + "_extracted")
        os.makedirs(extract_dir, exist_ok=True)
        
        try:
            with zipfile.ZipFile(zip_file_path, 'r') as zip_ref:
                # Get list of files in the zip
                file_list = [f for f in zip_ref.namelist() if not f.endswith('/')]
                
                # Extract each file
                for file_name in file_list:
                    # Extract the file
                    extracted_path = os.path.join(extract_dir, os.path.basename(file_name))
                    
                    # Use a unique filename to avoid collisions
                    unique_filename = str(uuid.uuid4()) + "_" + os.path.basename(file_name)
                    extracted_path = os.path.join(extract_dir, unique_filename)
                    
                    # Extract the file
                    with zip_ref.open(file_name) as source, open(extracted_path, 'wb') as target:
                        shutil.copyfileobj(source, target)
                    
                    # Create a FileStorage-like object for compatibility with the rest of the code
                    file_storage = FileStorage(
                        stream=open(extracted_path, 'rb'),
                        filename=os.path.basename(file_name),
                        name=os.path.basename(file_name)
                    )
                    
                    extracted_files.append((extracted_path, file_storage))
                    
            return extracted_files
            
        except Exception as e:
            logger.error(f"Error extracting ZIP file: {str(e)}")
            # Clean up the extraction directory if there was an error
            shutil.rmtree(extract_dir, ignore_errors=True)
            return []
    
    def extract_text(self, file_path, use_docling=False):
        """
        Extract text from various file formats, with option to use docling or fallback methods
        
        Args:
            file_path (str): Path to the file
            use_docling (bool): Whether to use Docling (True) or fallback methods (False)
            
        Returns:
            str: Extracted text content
        """
        # Determine file type based on extension
        file_extension = os.path.splitext(file_path)[1].lower()
        mime_type, _ = mimetypes.guess_type(file_path)
        
        # Try using docling for processing if enabled
        if use_docling and file_extension in self.extension_map:
            logger.info(f"Using Docling to process {file_extension} file")
            
            # Create a temporary working directory for Docling
            with tempfile.TemporaryDirectory() as temp_dir:
                try:
                    # Create a copy of the file in the temp directory to avoid permission issues
                    temp_file = os.path.join(temp_dir, os.path.basename(file_path))
                    shutil.copy2(file_path, temp_file)
                    
                    # Convert file path to Path object
                    path_obj = Path(temp_file)
                    
                    # Use docling to process document
                    result = self.converter.convert(source=path_obj)
                    
                    # Extract text - based on the example, we need to access document
                    if result and hasattr(result, 'document'):
                        # Get text using markdown export (ensures consistent formatting)
                        text = result.document.export_to_markdown()
                        
                        if text and text.strip():
                            logger.info(f"Successfully extracted text with Docling (text length: {len(text)})")
                            return text
                        
                        # Try getting text from raw content if markdown export is empty
                        if not text.strip() and hasattr(result.document, 'content'):
                            text = result.document.content
                            if text.strip():
                                logger.info(f"Successfully extracted raw content with Docling")
                                return text
                            
                        logger.warning(f"Docling returned empty text for {file_path}")
                    else:
                        logger.warning(f"Docling conversion failed to return document for {file_path}")
                        logger.debug(f"Result properties: {dir(result)}")
                except Exception as e:
                    logger.error(f"Docling conversion failed for {file_path}: {str(e)}")
        
        # If docling is disabled, fails, or not applicable, use fallback methods
        logger.info(f"Using fallback extraction for {file_path}")
        
        # Specialized handling based on file type
        
        # Handle Office documents
        if file_extension in ['.docx', '.doc']:
            try:
                import docx
                doc = docx.Document(file_path)
                text = "\n\n".join([paragraph.text for paragraph in doc.paragraphs])
                return text
            except Exception as e:
                logger.error(f"Error extracting Word document: {str(e)}")
                
        elif file_extension in ['.pptx', '.ppt']:
            try:
                from pptx import Presentation
                presentation = Presentation(file_path)
                text = ""
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                    text += "\n---\n"
                return text
            except Exception as e:
                logger.error(f"Error extracting PowerPoint: {str(e)}")
                
        elif file_extension in ['.xlsx', '.xls']:
            try:
                import pandas as pd
                sheets = pd.read_excel(file_path, sheet_name=None)
                text = ""
                for sheet_name, df in sheets.items():
                    text += f"Sheet: {sheet_name}\n"
                    text += df.to_string(index=False) + "\n\n"
                return text
            except Exception as e:
                logger.error(f"Error extracting Excel: {str(e)}")
        
        # Handle PDF files
        if file_extension == '.pdf' or mime_type == 'application/pdf':
            try:
                # Extract text from PDF
                reader = PdfReader(file_path)
                text = ""
                for i, page in enumerate(reader.pages):
                    text += f"--- Page {i+1} ---\n"
                    text += page.extract_text() + "\n\n"
                return text
            except Exception as e:
                logger.error(f"Error extracting PDF: {str(e)}")
            
        # Handle HTML files
        elif file_extension in ['.html', '.htm'] or mime_type == 'text/html':
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                soup = BeautifulSoup(content, 'html.parser')
                return soup.get_text()
            except Exception as e:
                logger.error(f"Error extracting HTML: {str(e)}")
            
        # Handle Markdown files
        elif file_extension == '.md' or mime_type == 'text/markdown':
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                html = markdown.markdown(content)
                soup = BeautifulSoup(html, 'html.parser')
                return soup.get_text()
            except Exception as e:
                logger.error(f"Error extracting Markdown: {str(e)}")
            
        # Handle text files
        elif file_extension in ['.txt', '.text'] or (mime_type and mime_type.startswith('text/')):
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            except Exception as e:
                logger.error(f"Error extracting text file: {str(e)}")
                
        # Attempt to read as raw text as last resort
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception as e:
            logger.error(f"Failed to read file as text: {str(e)}")
            return f"Could not extract text from {os.path.basename(file_path)} (format: {file_extension})"
    
    def chunk_text(self, text, chunk_size=1000, overlap=200, strategy="sentence"):
        """
        Split text into overlapping chunks using the specified strategy
        
        Args:
            text (str): Text to chunk
            chunk_size (int, optional): Size of each chunk in characters
            overlap (int, optional): Overlap between chunks in characters
            strategy (str, optional): Chunking strategy ('sentence', 'paragraph', 'fixed')
            
        Returns:
            list: List of text chunks
        """
        if not text:
            return []
        
        # Clean text (normalize whitespace but preserve paragraph breaks)
        if strategy == "paragraph":
            # Preserve paragraph breaks
            paragraphs = re.split(r'\n\s*\n', text)
            # Clean each paragraph
            paragraphs = [re.sub(r'\s+', ' ', p).strip() for p in paragraphs]
            # Filter out empty paragraphs
            paragraphs = [p for p in paragraphs if p]
        else:
            # For other strategies, normalize all whitespace
            text = re.sub(r'\s+', ' ', text).strip()
        
        chunks = []
        
        # Choose chunking strategy
        if strategy == "paragraph":
            chunks = self._chunk_by_paragraph(paragraphs, chunk_size, overlap)
        elif strategy == "sentence":
            chunks = self._chunk_by_sentence(text, chunk_size, overlap)
        elif strategy == "fixed":
            chunks = self._chunk_fixed_size(text, chunk_size, overlap)
        else:
            # Default to sentence-based chunking
            chunks = self._chunk_by_sentence(text, chunk_size, overlap)
        
        return chunks
    
    def _chunk_by_sentence(self, text, chunk_size, overlap):
        """Helper method for sentence-based chunking"""
        # Split text into sentences using improved regex
        # This matches sentence endings with punctuation followed by space or end of string
        sentences = re.split(r'(?<=[.!?])\s+|(?<=[.!?])$', text)
        sentences = [s for s in sentences if s.strip()]
        
        chunks = []
        current_chunk = ""
        
        for sentence in sentences:
            # If adding this sentence doesn't exceed chunk size, add it
            if len(current_chunk) + len(sentence) <= chunk_size:
                current_chunk += sentence + " "
            else:
                # Save current chunk if it's not empty
                if current_chunk:
                    chunks.append(current_chunk.strip())
                
                # Start new chunk with overlap if possible
                if overlap > 0 and current_chunk:
                    # Calculate how many characters to include in overlap
                    # Try to end at word boundaries
                    overlap_chars = min(overlap, len(current_chunk))
                    overlap_text = current_chunk[-overlap_chars:].lstrip()
                    
                    # If overlap doesn't start at a word boundary, adjust it
                    if overlap_chars > 0 and current_chunk[-overlap_chars] != ' ' and len(current_chunk) > overlap_chars:
                        # Find the nearest word boundary
                        space_pos = current_chunk[:-overlap_chars].rfind(' ')
                        if space_pos >= 0:
                            overlap_text = current_chunk[space_pos+1:]
                    
                    current_chunk = overlap_text + " " + sentence + " "
                else:
                    current_chunk = sentence + " "
        
        # Add the last chunk if it's not empty
        if current_chunk.strip():
            chunks.append(current_chunk.strip())
            
        return chunks
    
    def _chunk_by_paragraph(self, paragraphs, chunk_size, overlap):
        """Helper method for paragraph-based chunking"""
        chunks = []
        current_chunk = ""
        current_paragraphs = []
        
        for paragraph in paragraphs:
            # If adding this paragraph doesn't exceed chunk size, add it
            if len(current_chunk) + len(paragraph) + 2 <= chunk_size:  # +2 for \n\n
                if current_chunk:
                    current_chunk += "\n\n" + paragraph
                    current_paragraphs.append(paragraph)
                else:
                    current_chunk = paragraph
                    current_paragraphs.append(paragraph)
            else:
                # Save current chunk if it's not empty
                if current_chunk:
                    chunks.append(current_chunk)
                
                # Start new chunk with overlap if possible
                if overlap > 0 and current_paragraphs:
                    # Include previous paragraphs that fit within overlap
                    overlap_text = ""
                    overlap_paragraphs = []
                    
                    # Start from the end and add paragraphs until we exceed overlap
                    for p in reversed(current_paragraphs):
                        if len(p) + 2 + len(overlap_text) <= overlap:  # +2 for \n\n
                            if overlap_text:
                                overlap_text = p + "\n\n" + overlap_text
                            else:
                                overlap_text = p
                            overlap_paragraphs.insert(0, p)
                        else:
                            break
                    
                    current_chunk = overlap_text
                    if overlap_text:
                        current_chunk += "\n\n"
                    current_chunk += paragraph
                    current_paragraphs = overlap_paragraphs + [paragraph]
                else:
                    current_chunk = paragraph
                    current_paragraphs = [paragraph]
        
        # Add the last chunk if it's not empty
        if current_chunk:
            chunks.append(current_chunk)
            
        return chunks
    
    def _chunk_fixed_size(self, text, chunk_size, overlap):
        """Helper method for fixed-size chunking"""
        chunks = []
        
        # Simple character-based chunking with fixed size
        for i in range(0, len(text), chunk_size - overlap):
            if i > 0:
                # Start 'overlap' characters before to create overlap
                i = max(0, i - overlap)
            
            # Get chunk of appropriate size
            end = min(i + chunk_size, len(text))
            chunk = text[i:end]
            
            # Try to break at word boundaries if possible
            if end < len(text) and text[end] != ' ' and ' ' in chunk:
                # Find last space in chunk
                last_space = chunk.rfind(' ')
                if last_space > 0:
                    chunk = chunk[:last_space]
                    i = i + last_space
            
            # Only add non-empty chunks
            if chunk.strip():
                chunks.append(chunk.strip())
        
        return chunks
    
    def extract_enhanced_metadata(self, text, file_path):
        """
        Extract enhanced metadata from the document text
        
        Args:
            text (str): Document text
            file_path (str): Path to the document file
            
        Returns:
            dict: Enhanced metadata
        """
        metadata = {}
        
        # Basic document stats
        metadata['char_count'] = len(text)
        metadata['word_count'] = len(text.split())
        metadata['line_count'] = text.count('\n') + 1
        
        # Try to detect language (simple heuristic)
        # For a production system, use a proper language detector library
        common_english_words = {'the', 'and', 'of', 'to', 'in', 'a', 'is', 'that', 'for', 'it'}
        common_spanish_words = {'el', 'la', 'de', 'y', 'en', 'un', 'una', 'es', 'que', 'para'}
        common_french_words = {'le', 'la', 'les', 'de', 'et', 'un', 'une', 'est', 'que', 'pour'}
        
        words = set(word.lower() for word in re.findall(r'\b\w+\b', text.lower()))
        english_matches = len(words.intersection(common_english_words))
        spanish_matches = len(words.intersection(common_spanish_words))
        french_matches = len(words.intersection(common_french_words))
        
        if english_matches > spanish_matches and english_matches > french_matches:
            metadata['language'] = 'english'
        elif spanish_matches > english_matches and spanish_matches > french_matches:
            metadata['language'] = 'spanish'
        elif french_matches > english_matches and french_matches > spanish_matches:
            metadata['language'] = 'french'
        else:
            metadata['language'] = 'unknown'
        
        # Try to extract title
        lines = text.split('\n')
        if lines and lines[0].strip():
            # Assume first non-empty line might be a title if it's reasonably short
            first_line = lines[0].strip()
            if len(first_line) <= 200 and not first_line.endswith(('.', '?', '!')):
                metadata['potential_title'] = first_line
        
        # Extract potential headers
        headers = []
        for line in lines:
            line = line.strip()
            # Look for short lines that might be headers (not ending with sentence-ending punctuation)
            if line and len(line) <= 100 and not line.endswith(('.', '?', '!')):
                if re.match(r'^[A-Z0-9]', line):  # Starts with uppercase or number
                    headers.append(line)
        
        if headers:
            metadata['potential_headers'] = headers[:10]  # Limit to avoid oversized metadata
        
        # Attempt to detect document type based on content patterns
        if re.search(r'(copyright|all rights reserved|©)', text.lower()):
            metadata['contains_copyright'] = True
        
        if re.search(r'(confidential|internal use only|not for distribution)', text.lower()):
            metadata['confidential'] = True
        
        # Extract dates if present (simple regex for common formats)
        dates = re.findall(r'\b\d{1,2}[/-]\d{1,2}[/-]\d{2,4}\b|\b\d{4}[/-]\d{1,2}[/-]\d{1,2}\b', text)
        if dates:
            metadata['dates_found'] = dates[:5]  # Limit to first 5 dates
            
        return metadata
        
    def process_document(self, collection_id, file, title, use_docling=False, chunk_size=1000, 
                         chunk_overlap=200, chunk_strategy="sentence", extract_metadata="basic", file_path=None):
        """
        Process an uploaded document - extract text, chunk it, and store in Qdrant
        
        Args:
            collection_id (str): ID of the collection
            file: Uploaded file object
            title (str): Document title
            use_docling (bool): Whether to use Docling for document processing
            chunk_size (int): Size of each chunk in characters
            chunk_overlap (int): Overlap between chunks in characters
            chunk_strategy (str): Chunking strategy to use
            extract_metadata (str): Metadata extraction level (basic or enhanced)
            file_path (str, optional): If provided, use this path instead of saving the file
            
        Returns:
            Document: Created document
        """
        # Save file if not already saved (like from zip extraction)
        if file_path is None:
            file_path = self.save_uploaded_file(file)
        
        # Extract text using the specified processor
        text = self.extract_text(file_path, use_docling=use_docling)
        
        # Prepare base metadata
        metadata = {
            'filename': file.filename,
            'size': os.path.getsize(file_path),
            'chunk_count': 0,
            'processed_with': 'docling' if use_docling else 'fallback',
            'chunking': {
                'strategy': chunk_strategy,
                'size': chunk_size,
                'overlap': chunk_overlap
            }
        }
        
        # Add enhanced metadata if requested
        if extract_metadata == "enhanced":
            enhanced_metadata = self.extract_enhanced_metadata(text, file_path)
            metadata.update({'enhanced': enhanced_metadata})
        
        # Create document
        document = Document(
            collection_id=collection_id,
            title=title,
            content=text[:1000] + "..." if len(text) > 1000 else text,  # Store preview
            file_path=file_path,
            metadata=metadata
        )
        
        # Save document
        document.save()
        
        # Process document for Qdrant - chunk text and store
        chunks = self.chunk_text(text, chunk_size=chunk_size, 
                               overlap=chunk_overlap, strategy=chunk_strategy)
        
        # Get collection to determine embedding model to use
        collection = Collection.get(collection_id)
        
        # Ensure collection exists in Qdrant
        if not self.qdrant_service.collection_exists(collection_id):
            # Collection must exist and have embedding information
            if not collection:
                raise ValueError(f"Collection {collection_id} not found")
                
            if not collection.embedding_provider or not collection.embedding_model:
                raise ValueError(f"Collection {collection_id} does not have embedding model information")
                
            success = self.qdrant_service.create_collection(
                collection_name=collection_id, 
                vector_size=collection.embedding_dimensions,
                provider=collection.embedding_provider,
                model=collection.embedding_model
            )
            
            if not success:
                raise ValueError(f"Failed to create vector collection for {collection_id}")
        
        # Insert chunks into Qdrant
        for i, chunk in enumerate(chunks):
            # We'll use a simple index as reference and let the Qdrant service handle UUID generation
            chunk_metadata = {
                'document_id': document.id,
                'document_title': title,
                'chunk_index': i,
                'original_chunk_id': f"{document.id}_{i}",
                'content': chunk,
                'file_path': document.file_path,
                'chunk_strategy': chunk_strategy,
                'chunk_size': chunk_size,
                'chunk_overlap': chunk_overlap
            }
            
            # Add page number approximation
            if "--- Page " in chunk:
                # Try to extract actual page number from text
                page_match = re.search(r"--- Page (\d+) ---", chunk)
                if page_match:
                    chunk_metadata['page_number'] = int(page_match.group(1))
                else:
                    chunk_metadata['page_number'] = i + 1  # Simple approximation
            else:
                chunk_metadata['page_number'] = i + 1  # Simple approximation
            
            # We'll use a simple reference ID and let Qdrant service generate a proper UUID
            ref_id = f"{document.id}_{i}"
            
            # Use collection's embedding model if defined
            if collection and collection.embedding_provider and collection.embedding_model:
                self.qdrant_service.insert_document(
                    collection_name=collection_id,
                    document_id=ref_id,
                    text=chunk,
                    metadata=chunk_metadata,
                    provider=collection.embedding_provider,
                    model=collection.embedding_model
                )
            else:
                self.qdrant_service.insert_document(
                    collection_name=collection_id,
                    document_id=ref_id,
                    text=chunk,
                    metadata=chunk_metadata
                )
        
        # Update document metadata
        document.metadata['chunk_count'] = len(chunks)
        document.save()
        
        return document
    
    def search_documents(self, collection_id, query, limit=5):
        """
        Search for documents in a collection
        
        Args:
            collection_id (str): ID of the collection
            query (str): Search query
            limit (int, optional): Maximum number of results
            
        Returns:
            dict: Search results with text chunks and document info
        """
        if not self.qdrant_service.collection_exists(collection_id):
            return {
                'query': query,
                'results': [],
                'message': 'Collection does not exist'
            }
        
        # Get collection to determine embedding model to use
        collection = Collection.get(collection_id)
        
        # Search Qdrant with the collection's embedding model if defined
        if collection and collection.embedding_provider and collection.embedding_model:
            results = self.qdrant_service.search(
                collection_id, 
                query, 
                limit=limit,
                provider=collection.embedding_provider,
                model=collection.embedding_model
            )
        else:
            results = self.qdrant_service.search(collection_id, query, limit=limit)
        
        # Format search results
        formatted_results = []
        for result in results:
            payload = result['payload']
            formatted_results.append({
                'score': result['score'],
                'document_id': payload.get('document_id'),
                'document_title': payload.get('document_title'),
                'content': payload.get('content', ''),
                'file_path': payload.get('file_path', ''),
                'page_number': payload.get('page_number', 1)
            })
        
        return {
            'query': query,
            'results': formatted_results
        }